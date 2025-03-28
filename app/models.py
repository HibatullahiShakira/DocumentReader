import pdfplumber
import pptx
import re
from nltk.sentiment import SentimentIntensityAnalyzer
from nltk.tokenize import word_tokenize, sent_tokenize
from nltk.corpus import stopwords
from nltk.tag import pos_tag
import nltk
from collections import Counter
from flask_sqlalchemy import SQLAlchemy
from datetime import datetime, UTC

nltk.download('punkt')
nltk.download('punkt_tab')
nltk.download('averaged_perceptron_tagger')
nltk.download('averaged_perceptron_tagger_eng')
nltk.download('stopwords')
nltk.download('vader_lexicon')

db = SQLAlchemy()


class PitchDeckParser:
    def __init__(self, sia=None):
        self.sia = sia if sia else SentimentIntensityAnalyzer()
        self.stop_words = set(stopwords.words('english'))

    def parse_pdf(self, file_path):
        try:
            with pdfplumber.open(file_path) as pdf:
                content = ""
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        page_text = re.sub(r'\s+', ' ', page_text.strip())
                        content += page_text + "\n"
                content = content.rstrip()
                return content, len(pdf.pages)
        except Exception as e:
            print(f"PDF parsing error: {e}")
            raise

    def parse_pptx(self, file_path):
        try:
            prs = pptx.Presentation(file_path)
            content = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
            return content, len(prs.slides)
        except Exception as e:
            print(f"PPTX parsing error: {e}")
            raise

    def detect_document_type(self, text):
        text_lower = text.lower()
        pitch_deck_patterns = [
            r'our problem is\s+[^.\n]+',
            r'our solution is\s+[^.\n]+',
            r'the market is\s+[^.\n]+',
            r'problem is\s+[^.\n]+',
            r'solution is\s+[^.\n]+'
        ]
        for pattern in pitch_deck_patterns:
            if re.search(pattern, text_lower):
                print(f"Found pitch deck pattern: {pattern}")
                return 'pitch_deck'

        personal_sections = ['objective', 'summary', 'profile']
        other_sections = ['experience', 'education', 'skills', 'certifications']
        has_personal_section = False
        has_other_section = False

        for keyword in personal_sections:
            if re.search(rf'(?i)\b{keyword}\b(?:\s*$|\s+.*$)', text_lower, re.MULTILINE):
                print(f"Found personal section: {keyword}")
                has_personal_section = True
                break

        for keyword in other_sections:
            if re.search(rf'(?i)\b{keyword}(?:\s*(?:&.*)?(?:\s*$|\s+.*$))', text_lower, re.MULTILINE):
                print(f"Found other section: {keyword}")
                has_other_section = True
                break

        print(f"has_personal_section: {has_personal_section}, has_other_section: {has_other_section}")
        if has_personal_section and has_other_section:
            return 'resume'

        return 'generic'

    def extract_section(self, lines, start_keyword, max_lines=10):
        for i, line in enumerate(lines):
            line_lower = line.lower().strip()
            if re.search(rf'(?i)\b{start_keyword}(?:\s*(?:&.*)?(?:\s*:.*|\s+.*)?$)', line_lower):
                section_lines = [line.strip()]
                for j in range(1, max_lines + 1):
                    if i + j < len(lines):
                        next_line = lines[i + j].strip()
                        next_line_lower = next_line.lower()
                        if not next_line or next_line_lower.startswith(('objective', 'summary', 'profile', 'experience',
                                                                        'skills', 'education', 'certifications')):
                            break
                        section_lines.append(next_line)
                    else:
                        break
                return re.sub(r'\s+', ' ', ' '.join(section_lines))
            if start_keyword == 'experience':
                if re.search(
                        r'(?i)[a-z\s&-]+(?:intern|engineer|developer|analyst|manager|specialist)(?:\s*:.*|\s+.*)?$',
                        line_lower):
                    section_lines = [line.strip()]
                    for j in range(1, max_lines + 1):
                        if i + j < len(lines):
                            next_line = lines[i + j].strip()
                            next_line_lower = next_line.lower()
                            if not next_line or next_line_lower.startswith(('objective', 'summary', 'profile',
                                                                            'experience', 'skills', 'education',
                                                                            'certifications')):
                                break
                            section_lines.append(next_line)
                        else:
                            break
                    return re.sub(r'\s+', ' ', ' '.join(section_lines))
        return None

    def extract_key_phrases(self, text, top_n=5):
        words = word_tokenize(text.lower())
        words = [word for word in words if word.isalnum() or '-' in word]
        words = [word for word in words if word not in self.stop_words]
        tagged_words = pos_tag(words)
        phrases = []
        current_phrase = []
        phrase_positions = {}

        for word, tag in tagged_words:
            if tag.startswith(('NN', 'JJ')) or '-' in word:
                current_phrase.append(word)
            else:
                if current_phrase:
                    for length in range(2, 4):
                        for start in range(len(current_phrase) - length + 1):
                            if start + length <= len(current_phrase):
                                phrase = ' '.join(current_phrase[start:start + length])
                                if phrase not in phrase_positions:
                                    phrase_positions[phrase] = words.index(current_phrase[start])
                                phrases.append(phrase)
                current_phrase = []

        if current_phrase:
            for length in range(2, 4):
                for start in range(len(current_phrase) - length + 1):
                    if start + length <= len(current_phrase):
                        phrase = ' '.join(current_phrase[start:start + length])
                        if phrase not in phrase_positions:
                            phrase_positions[phrase] = words.index(current_phrase[start])
                        phrases.append(phrase)

        phrase_counts = Counter(phrases)
        sorted_phrases = sorted(
            phrase_counts.items(),
            key=lambda x: (-x[1], -len(x[0].split()), phrase_positions[x[0]])
        )
        key_phrases = [phrase for phrase, count in sorted_phrases[:top_n] if len(phrase.split()) > 1]
        return key_phrases if key_phrases else ["No key phrases identified"]

    def extract_summary(self, text, max_sentences=3):
        sentences = sent_tokenize(text.strip())
        if not sentences:
            return "No content to summarize."

        key_phrases = self.extract_key_phrases(text, top_n=5)
        if not key_phrases or key_phrases == ["No key phrases identified"]:
            return ' '.join(sentences[:max_sentences]).strip()

        sentence_scores = []
        for sentence in sentences:
            score = sum(1 for phrase in key_phrases if phrase.lower() in sentence.lower())
            sentence_scores.append((sentence, score))

        sentence_scores.sort(key=lambda x: (-x[1], sentences.index(x[0])))
        top_sentences = [sentence for sentence, score in sentence_scores[:max_sentences]]
        return ' '.join(top_sentences).strip()

    def analyze_content(self, text):
        info = {
            'upload_date': datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            'word_count': len(text.split()),
            'char_count': len(text.replace('\n', ''))
        }
        sentiment = self.sia.polarity_scores(text)
        info['sentiment_score'] = sentiment['compound']
        info['sentiment_type'] = 'Positive' if sentiment['compound'] > 0.05 else \
            'Negative' if sentiment['compound'] < -0.05 else 'Neutral'

        doc_type = self.detect_document_type(text.lower())
        info['document_type'] = doc_type
        lines = text.lower().split('\n')

        if doc_type == 'pitch_deck':
            for line in lines:
                if 'problem' in line:
                    info['problem'] = line.strip()
                elif 'solution' in line:
                    info['solution'] = line.strip()
                elif 'market' in line:
                    info['market'] = line.strip()
        elif doc_type == 'resume':
            info['problem'] = self.extract_section(lines, 'objective') or \
                              self.extract_section(lines, 'summary') or \
                              self.extract_section(lines, 'profile')
            info['experience'] = self.extract_section(lines, 'experience')
            info['skills'] = self.extract_section(lines, 'skills')
        else:
            info['key_phrases'] = self.extract_key_phrases(text, top_n=5)
            info['summary'] = self.extract_summary(text)

        if 'problem' not in info or not info['problem']:
            info['problem'] = self.extract_summary(text)

        return info


class PitchDeck(db.Model):
    __tablename__ = 'pitch_decks'

    id = db.Column(db.Integer, primary_key=True)
    filename = db.Column(db.String(255), nullable=False)
    upload_date = db.Column(db.DateTime, default=lambda: datetime.now(UTC))
    content = db.Column(db.Text)
    slide_count = db.Column(db.Integer)
    status = db.Column(db.String(50))
    word_count = db.Column(db.Integer)
    char_count = db.Column(db.Integer)
    sentiment_score = db.Column(db.Float)
    sentiment_type = db.Column(db.String(50))
    document_type = db.Column(db.String(50))
    problem = db.Column(db.Text)
    solution = db.Column(db.Text)
    market = db.Column(db.Text)
    experience = db.Column(db.Text)
    skills = db.Column(db.Text)
    summary = db.Column(db.Text)
    key_phrases = db.Column(db.Text)

    def __init__(self, filename, content, slide_count, analysis, status='processed'):
        self.filename = filename
        self.content = content
        self.slide_count = slide_count
        self.status = status
        self.analysis = analysis
        self.word_count = analysis['word_count']
        self.char_count = analysis['char_count']
        self.sentiment_score = analysis['sentiment_score']
        self.sentiment_type = analysis['sentiment_type']
        self.document_type = analysis.get('document_type')
        self.problem = analysis.get('problem')
        self.solution = analysis.get('solution')
        self.market = analysis.get('market')
        self.experience = analysis.get('experience')
        self.skills = analysis.get('skills')
        self.summary = analysis.get('summary')
        self.key_phrases = ', '.join(analysis.get('key_phrases', [])) if analysis.get('key_phrases') else None
        upload_date_str = analysis.get('upload_date')
        if upload_date_str:
            self.upload_date = datetime.strptime(upload_date_str, "%Y-%m-%d %H:%M:%S")

    def save(self, redis_client):
        try:
            db.session.add(self)
            db.session.commit()
            redis_client.delete('dashboard_data')
            return self.id
        except Exception as e:
            db.session.rollback()
            print(f"Database storage error: {e}")
            raise
