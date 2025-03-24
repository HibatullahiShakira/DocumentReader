from flask import Flask
import os
import PyPDF2
import pptx
from datetime import datetime
import psycopg2
from psycopg2 import Error
import redis
import json
from nltk.sentiment import SentimentIntensityAnalyzer
import nltk
nltk.download('vader_lexicon')

from app.routes import init_routes

def create_app():
    app = Flask(__name__)

    UPLOAD_FOLDER = 'uploads'
    MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB
    app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
    app.config['MAX_CONTENT_LENGTH'] = MAX_FILE_SIZE

    DB_CONFIG = {
        'dbname': os.getenv('POSTGRES_DB', 'pitch_decks'),
        'user': os.getenv('POSTGRES_USER', 'postgres'),
        'password': os.getenv('POSTGRES_PASSWORD', 'password'),
        'host': os.getenv('POSTGRES_HOST', 'db'),
        'port': os.getenv('POSTGRES_PORT', '5432')
    }

    try:
        redis_client = redis.Redis(
            host=os.getenv('REDIS_HOST', 'redis'),
            port=int(os.getenv('REDIS_PORT', '6379')),
            db=0,
            decode_responses=True
        )
        redis_client.ping()
        print("Successfully connected to Redis")
    except redis.ConnectionError as e:
        print(f"Failed to connect to Redis: {e}")
        raise

    sia = SentimentIntensityAnalyzer()

    print("Starting Flask application")

    def init_db():
        print("Initializing database")
        try:
            conn = psycopg2.connect(**DB_CONFIG)
            cur = conn.cursor()
            cur.execute("""
                CREATE TABLE IF NOT EXISTS pitch_decks (
                    id SERIAL PRIMARY KEY,
                    filename VARCHAR(255) NOT NULL,
                    upload_date TIMESTAMP,
                    content TEXT,
                    slide_count INTEGER,
                    status VARCHAR(50),
                    word_count INTEGER,
                    char_count INTEGER,
                    sentiment_score FLOAT,
                    sentiment_type VARCHAR(50),
                    problem TEXT,
                    solution TEXT,
                    market TEXT
                )
            """)
            conn.commit()
            cur.close()
            conn.close()
            print("Database initialized successfully")
        except Error as e:
            print(f"Database initialization error: {e}")
            raise

    class PitchDeckParser:
        def parse_pdf(self, file_path):
            try:
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    content = ""
                    for page in pdf_reader.pages:
                        content += page.extract_text() + "\n"
                    return content, len(pdf_reader.pages)
            except Exception as e:
                print(f"PDF parsing error: {e}")
                raise

        def parse_pptx(self, file_path):
            try:
                prs = pptx.Presentation(file_path)
                content = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            content += shape.text + "\n"
                    return content, len(prs.slides)
            except Exception as e:
                print(f"PPTX parsing error: {e}")
                raise

        def analyze_content(self, text):
            info = {
                'upload_date': datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                'word_count': len(text.split()),
                'char_count': len(text)
            }
            sentiment = sia.polarity_scores(text)
            info['sentiment_score'] = sentiment['compound']
            info['sentiment_type'] = 'Positive' if sentiment['compound'] > 0.05 else \
                'Negative' if sentiment['compound'] < -0.05 else 'Neutral'
            lines = text.lower().split('\n')
            for line in lines:
                if 'problem' in line:
                    info['problem'] = line.strip()
                elif 'solution' in line:
                    info['solution'] = line.strip()
                elif 'market' in line:
                    info['market'] = line.strip()
            return info

        def store_data(self, filename, content, slide_count, analysis):
            try:
                conn = psycopg2.connect(**DB_CONFIG)
                cur = conn.cursor()
                cur.execute("""
                    INSERT INTO pitch_decks (
                        filename, upload_date, content, slide_count, status,
                        word_count, char_count, sentiment_score, sentiment_type,
                        problem, solution, market
                    )
                    VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s)
                    RETURNING id
                """, (
                    filename,
                    datetime.now(),
                    content,
                    slide_count,
                    'processed',
                    analysis['word_count'],
                    analysis['char_count'],
                    analysis['sentiment_score'],
                    analysis['sentiment_type'],
                    analysis.get('problem'),
                    analysis.get('solution'),
                    analysis.get('market')
                ))
                deck_id = cur.fetchone()[0]
                conn.commit()
                cur.close()
                conn.close()
                redis_client.delete('dashboard_data')
                return deck_id
            except Error as e:
                print(f"Database storage error: {e}")
                raise

    parser = PitchDeckParser()
    init_routes(app, redis_client, parser, DB_CONFIG)

    print("Creating upload folder")
    os.makedirs(UPLOAD_FOLDER, exist_ok=True)
    print("Calling init_db")
    init_db()

    return app

app = create_app()

if __name__ == '__main__':
    print("Starting Flask server")
    app.run(host='0.0.0.0', port=5000)