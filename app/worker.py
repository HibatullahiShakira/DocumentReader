import logger
import redis
import json
import os
import logging
from datetime import datetime
import psycopg2
from psycopg2 import Error
import PyPDF2
import pptx
from nltk.sentiment import SentimentIntensityAnalyzer
import nltk

# Redis configuration
try:
    redis_client = redis.Redis(
        host=os.getenv('REDIS_HOST', 'redis'),
        port=int(os.getenv('REDIS_PORT', '6379')),
        db=0,
        decode_responses=True
    )
    redis_client.ping()  # Test the connection
    logger.info("Successfully connected to Redis")
except redis.ConnectionError as e:
    logger.error(f"Failed to connect to Redis: {e}")
    raise

# Database configuration
DB_CONFIG = {
    'dbname': os.getenv('POSTGRES_DB', 'pitch_decks'),
    'user': os.getenv('POSTGRES_USER', 'postgres'),
    'password': os.getenv('POSTGRES_PASSWORD', 'password'),
    'host': os.getenv('POSTGRES_HOST', 'db'),
    'port': os.getenv('POSTGRES_PORT', '5432')
}

# Initialize sentiment analyzer
sia = SentimentIntensityAnalyzer()

# Setup logging
logging.basicConfig(level=logging.DEBUG)  # Set to DEBUG for more detail
logger = logging.getLogger(__name__)
logger.debug("Starting worker")

class PitchDeckParser:
    def parse_pdf(self, file_path):
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                content = ""
                for page in pdf_reader.pages:
                    content += page.extract_text() + "\n"
                return content, len(pdf_reader.pages)
        except Exception as e:
            logger.error(f"PDF parsing error: {e}")
            raise

    def parse_pptx(self, file_path):
        try:
            prs = pptx.Presentation(file_path)
            content = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
            return content, len(prs.slides)
        except Exception as e:
            logger.error(f"PPTX parsing error: {e}")
            raise

    def analyze_content(self, text):
        info = {
            'upload_date': datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            'word_count': len(text.split()),
            'char_count': len(text)
        }

        sentiment = sia.polarity_scores(text)
        info['sentiment_score'] = sentiment['compound']
        info['sentiment_type'] = 'Positive' if sentiment['compound'] > 0.05 else \
            'Negative' if sentiment['compound'] < -0.05 else 'Neutral'

        lines = text.lower().split('\n')
        for line in lines:
            if 'problem' in line:
                info['problem'] = line.strip()
            elif 'solution' in line:
                info['solution'] = line.strip()
            elif 'market' in line:
                info['market'] = line.strip()

        return info

    def store_data(self, filename, content, slide_count, analysis):
        try:
            conn = psycopg2.connect(**DB_CONFIG)
            cur = conn.cursor()
            cur.execute("""
                INSERT INTO pitch_decks (
                    filename, upload_date, content, slide_count, status,
                    word_count, char_count, sentiment_score, sentiment_type,
                    problem, solution, market
                )
                VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s)
                RETURNING id
            """, (
                filename,
                datetime.now(),
                content,
                slide_count,
                'processed',
                analysis['word_count'],
                analysis['char_count'],
                analysis['sentiment_score'],
                analysis['sentiment_type'],
                analysis.get('problem'),
                analysis.get('solution'),
                analysis.get('market')
            ))
            deck_id = cur.fetchone()[0]
            conn.commit()
            cur.close()
            conn.close()

            # Clear cache after new data
            redis_client.delete('dashboard_data')
            return deck_id
        except Error as e:
            logger.error(f"Database storage error: {e}")
            raise

def process_queue():
    parser = PitchDeckParser()
    while True:
        try:
            logger.debug("Waiting for job in processing queue")
            job_data = redis_client.brpop('processing_queue', timeout=5)
            if job_data:
                _, job_json = job_data
                job = json.loads(job_json)
                file_path = job['file_path']
                filename = job['filename']

                # Process file
                logger.debug(f"Processing file: {filename}")
                if filename.endswith('.pdf'):
                    content, slide_count = parser.parse_pdf(file_path)
                else:
                    content, slide_count = parser.parse_pptx(file_path)

                # Analyze content
                logger.debug("Analyzing content")
                analysis = parser.analyze_content(content)

                # Store in database
                logger.debug("Storing data in database")
                parser.store_data(filename, content, slide_count, analysis)

                # Clean up
                logger.debug("Cleaning up temporary file")
                if os.path.exists(file_path):
                    os.remove(file_path)
                logger.info(f"Processed file: {filename}")
        except Exception as e:
            logger.error(f"Queue processing error: {e}")

if __name__ == '__main__':
    logger.debug("Starting queue processing")
    process_queue()