from app import celery, db
from app.models import File, Slide
import pdfplumber
from pptx import Presentation
import os


@celery.task(bind=True)
def process_file(self, file_id, filepath):
    try:
        file = File.query.get(file_id)
        if not file:
            raise ValueError("File not found")
        file.status = 'PROCESSING'
        db.session.commit()

        slides_data = []
        if filepath.endswith('.pdf'):
            try:
                with pdfplumber.open(filepath) as pdf:
                    metadata = pdf.metadata or {}
                    metadata['page_count'] = len(pdf.pages)
                    for i, page in enumerate(pdf.pages):
                        text = page.extract_text() or ''
                        lines = text.split('\n')
                        title = lines[0].strip() if lines else f"Slide {i + 1}"
                        slides_data.append({
                            'slide_number': i + 1,
                            'title': title,
                            'content': text,
                            'metadata': metadata
                        })
            except Exception as e:
                raise ValueError(f"Failed to process PDF: {str(e)}")
        elif filepath.endswith('.pptx'):
            try:
                prs = Presentation(filepath)
                metadata = {
                    'title': prs.core_properties.title,
                    'author': prs.core_properties.author,
                    'created': str(prs.core_properties.created),
                    'slide_count': len(prs.slides)
                }
                for i, slide in enumerate(prs.slides):
                    title = None
                    text = ''
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if not title and run.text.strip():
                                        title = run.text.strip()
                                    text += run.text + '\n'
                    slides_data.append({
                        'slide_number': i + 1,
                        'title': title or f"Slide {i + 1}",
                        'content': text.strip(),
                        'metadata': metadata
                    })
            except Exception as e:
                raise ValueError(f"Failed to process PPTX: {str(e)}")
        else:
            raise ValueError("Unsupported file format")

        for slide_data in slides_data:
            slide = Slide(
                file_id=file.id,
                slide_number=slide_data['slide_number'],
                title=slide_data['title'],
                content=slide_data['content'],
                metadata=slide_data['metadata']
            )
            db.session.add(slide)
        file.status = 'SUCCESS'
        db.session.commit()

        os.remove(filepath)

        return {'status': 'SUCCESS', 'slides': slides_data}
    except Exception as e:
        if 'file' in locals() and file:
            file.status = 'FAILURE'
            db.session.commit()
        raise e
